import streamlit as st
import base64
import os
import asyncio
import groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Set page config
st.set_page_config(page_title="GenPPT", layout="wide", page_icon="📊", initial_sidebar_state="collapsed")

# Load custom CSS
with open("style.css") as f:
    st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

# Initialize session state for navigation
if 'current_page' not in st.session_state:
    st.session_state.current_page = 'intro'

# Initialize Groq client with API key from .env
groq_api_key = os.getenv('GROQ_API_KEY')

if not groq_api_key:
    st.error("GROQ_API_KEY is not set. Please set it in your .env file.")
    st.stop()

client = groq.AsyncClient(api_key=groq_api_key)

async def generate_slide_titles(topic):
    prompt = f"Generate 10 slides for the topic '{topic}'."
    response = await client.chat.completions.create(
        model="llama-3.1-70b-versatile",
        messages=[
            {"role": "system", "content": "You are an expert in creating presentations. Your task is to generate slide titles by gathering comprehensive and accurate data"},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.split("\n")

async def generate_slide_content(slide_title):
    prompt = f"Generate detailed content for the slide: '{slide_title}'. Include a brief introduction, 3-4 key points with short descriptions, and a conclusion. Format the content with appropriate bullet points and sub-points."
    response = await client.chat.completions.create(
        model="llama-3.1-70b-versatile",
        messages=[
            {"role": "system", "content": "Your research should include:Industry trends: Identify current trends and benchmarks.Statistical analysis: Collect relevant statistics, percentages, and numerical data to support arguments.Case studies: Investigate real-world examples to make the presentation more relatable and credible.Comparative analysis: Look for comparisons between industry leaders, markets, or products, presenting pros and cons.Visual data: Ensure you gather suitable data for charts (bar, line, pie) and diagrams to make the content visually engaging.Organize the information into logical sections, prioritizing clarity and depth of analysis, ensuring no detail is overlooked. All research must be referenced and fact-checked, maintaining accuracy.."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

def create_ppt(topic, slide_titles, slide_contents):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "Generated Presentation"

    # Content slides
    for title, content in zip(slide_titles, slide_contents):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Set slide title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Add content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = content
        
        # Format content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            if paragraph.level == 0:
                paragraph.font.bold = True
    
    ppt_path = f"{topic}_presentation.pptx"
    prs.save(ppt_path)
    return ppt_path

async def generate_ppt(input_text):
    try:
        slide_titles = await generate_slide_titles(input_text)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        slide_contents = await asyncio.gather(*[generate_slide_content(title) for title in filtered_slide_titles])
        
        ppt_path = create_ppt(input_text, filtered_slide_titles, slide_contents)
        
        return {"slide_titles": filtered_slide_titles, "slide_contents": slide_contents, "ppt_path": ppt_path}
    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None

# Navigation
def nav_bar():
    col1, col2, col3, col4 = st.columns([1,1,1,3])
    with col1:
        if st.button("GenPPT", key="nav_intro", use_container_width=True):
            st.session_state.current_page = 'intro'
    with col2:
        if st.button("GENERATE", key="nav_generate", use_container_width=True):
            st.session_state.current_page = 'generate'
    with col3:
        if st.button("WHY?", key="nav_about", use_container_width=True):
            st.session_state.current_page = 'about'

# Pages
def intro_page():
    st.markdown('<h1 class="main-heading">Welcome to GenPPT</h1>', unsafe_allow_html=True)
    st.markdown("""
    <div class="intro-content-3d">
        <div class="intro-content">
            <h2 class="sub-heading">Transform Your Ideas into Stunning Presentations</h2>
            <p>GenPPT harnesses the power of Large Language Models to create professional PowerPoint presentations in seconds. Just enter your thoughts and have your presentation ready on the go.</p>
            <h3>Guidelines for best results:</h3>
            <ol>
                <li>Enter your presentation topic.</li>
                <li>Make sure to be specific and clear, add references if you can.</li>
                <li>Click 'Generate Presentation'.</li>
                <li>Review the generated slides.</li>
                <li>Download your ready-to-use PowerPoint file</li>
            </ol>
        </div>
    </div>
    """, unsafe_allow_html=True)

def generate_page():
    st.markdown('<h1 class="generate-heading">Generate Your Presentation</h1>', unsafe_allow_html=True)

    topic = st.text_input("Enter the topic for your presentation:")
    generate_button = st.button("Generate Presentation", key="generate_btn")

    if generate_button and topic:
        with st.spinner("Generating your presentation...This may take a few seconds."):
            try:
                result = asyncio.run(generate_ppt(topic))
                
                if result:
                    st.success("Presentation generated successfully!")
                    
                    # Display PPT preview
                    st.subheader("Presentation Preview")
                    for title, content in zip(result["slide_titles"], result["slide_contents"]):
                        with st.expander(title):
                            st.markdown(content)
                    
                    # Download button
                    st.markdown(get_ppt_download_link(result["ppt_path"]), unsafe_allow_html=True)
                else:
                    st.error("An error occurred while generating the presentation.")
            except Exception as e:
                st.error(f"An unexpected error occurred: {str(e)}")

def about_page():
    st.markdown('<h1 class="about-heading">About GenPPT</h1>', unsafe_allow_html=True)
    st.markdown("""
    <div class="about-content">
        <div class="engineer-message">
            <p>GenPPT aims to streamline the process by generating PowerPoint presentations from scratch.</p>
            <p>Instead of spending hours building slides, use this tool to brainstorm and organize your ideas efficiently.</p>
            <p>You can either copy the generated content into your file for further editing or download the complete PPT to customize it as needed.</p>
            <p>Easy!\n\n
            
            ggengineerco@gmail.com
            
    </div>
   
    """, unsafe_allow_html=True)

def get_ppt_download_link(ppt_path):
    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(ppt_path)}" class="download-btn">Download PowerPoint</a>'

def main():
    nav_bar()
    
    if st.session_state.current_page == 'intro':
        intro_page()
    elif st.session_state.current_page == 'generate':
        generate_page()
    elif st.session_state.current_page == 'about':
        about_page()

if __name__ == "__main__":
    main()
