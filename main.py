import os
from dotenv import load_dotenv
import groq
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
import asyncio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

app = FastAPI()

# Initialize Groq client with API key from .env
client = groq.AsyncClient(api_key=os.getenv("GROQ_API_KEY"))

class TextInput(BaseModel):
    text: str

async def generate_slide_titles(topic):
    prompt = f"Generate 10 slides for the topic '{topic}'."
    response = await client.chat.completions.create(
        model="llama-3.1-70b-versatile",
        messages=[
            {"role": "system", "content": "You are an expert in creating presentations. Your task is to generate slide titles by gathering comprehensive and accurate data"},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.split("\n")

async def generate_slide_content(slide_title):
    prompt = f"Generate detailed content for the slide: '{slide_title}'. Include a brief introduction, 3-4 key points with short descriptions, and a conclusion. Format the content with appropriate bullet points and sub-points."
    response = await client.chat.completions.create(
        model="llama-3.1-70b-versatile",
        messages=[
            {"role": "system", "content": "Your research should include:Industry trends: Identify current trends and benchmarks.Statistical analysis: Collect relevant statistics, percentages, and numerical data to support arguments.Case studies: Investigate real-world examples to make the presentation more relatable and credible.Comparative analysis: Look for comparisons between industry leaders, markets, or products, presenting pros and cons.Visual data: Ensure you gather suitable data for charts (bar, line, pie) and diagrams to make the content visually engaging.Organize the information into logical sections, prioritizing clarity and depth of analysis, ensuring no detail is overlooked. All research must be referenced and fact-checked, maintaining accuracy.."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

def create_ppt(topic, slide_titles, slide_contents):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "Generated Presentation"

    # Content slides
    for title, content in zip(slide_titles, slide_contents):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Set slide title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Add content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = content
        
        # Format content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            if paragraph.level == 0:
                paragraph.font.bold = True
    
    prs.save(f"generated_ppt/{topic}_presentation.pptx")
    return f"generated_ppt/{topic}_presentation.pptx"

@app.post("/generate_ppt")
async def generate_ppt(input: TextInput):
    try:
        slide_titles = await generate_slide_titles(input.text)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        slide_contents = await asyncio.gather(*[generate_slide_content(title) for title in filtered_slide_titles])
        
        ppt_path = create_ppt(input.text, filtered_slide_titles, slide_contents)
        
        return {"slide_titles": filtered_slide_titles, "slide_contents": slide_contents, "ppt_path": ppt_path}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

